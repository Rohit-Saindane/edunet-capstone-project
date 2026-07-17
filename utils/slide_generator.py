import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def apply_background(slide, prs):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)

def create_text_box(slide, left, top, width, height, margin=0):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    return tf

def add_slide_header(slide, title_text, category_text="EDUAI | SDG 4"):
    tf_cat = create_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.3))
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = "Arial"
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = RGBColor(139, 92, 246)
    tf_title = create_text_box(slide, Inches(0.8), Inches(0.6), Inches(10), Inches(0.8))
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Arial"
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 255, 255)

def generate_presentation(output_path: str = "report/presentation.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]
    
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide, prs)
    tf_sub = create_text_box(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(0.5))
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "UNITED NATIONS SDG 4 - QUALITY EDUCATION EXHIBITION"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(12)
    p_sub.font.bold = True
    p_sub.font.color.rgb = RGBColor(96, 165, 250)
    tf_title = create_text_box(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(1.5))
    p_title = tf_title.paragraphs[0]
    p_title.text = "EduAI: Personalized Learning Assistant"
    p_title.font.name = "Arial"
    p_title.font.size = Pt(46)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 255, 255)
    tf_desc = create_text_box(slide, Inches(1.5), Inches(4.3), Inches(9), Inches(1.0))
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = "A smart system leveraging Retrieval-Augmented Generation (RAG) and SQLite logs to provide individualized summaries, quizzes, and study planning instead of acting as a generic chatbot."
    p_desc.font.name = "Arial"
    p_desc.font.size = Pt(16)
    p_desc.font.color.rgb = RGBColor(203, 213, 225)
    tf_pres = create_text_box(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.8))
    p_pres = tf_pres.paragraphs[0]
    p_pres.text = "Developed for the SDG 4 College Exhibition | AI Solution Track"
    p_pres.font.name = "Arial"
    p_pres.font.size = Pt(12)
    p_pres.font.italic = True
    p_pres.font.color.rgb = RGBColor(148, 163, 184)
    
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide, prs)
    add_slide_header(slide, "UN SDG 4: Quality Education")
    tf_body = create_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
    p1 = tf_body.paragraphs[0]
    p1.text = "The Global Goal:"
    p1.font.bold = True
    p1.font.size = Pt(20)
    p1.font.color.rgb = RGBColor(96, 165, 250)
    p2 = tf_body.add_paragraph()
    p2.text = "UN SDG 4 aims to 'ensure inclusive and equitable quality education and promote lifelong learning opportunities for all' by 2030."
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(226, 232, 240)
    p2.space_after = Pt(20)
    p3 = tf_body.add_paragraph()
    p3.text = "Key Challenges We Address:"
    p3.font.bold = True
    p3.font.size = Pt(20)
    p3.font.color.rgb = RGBColor(96, 165, 250)
    bullets = [
        "Inaccessibility of premium tutoring services for lower-income students.",
        "The 'one-size-fits-all' curriculum which ignores individual student paces and learning blocks.",
        "Inefficient study tracking that fails to highlight conceptual weaknesses.",
        "Lack of custom-tailored feedback on assignments and test attempts."
    ]
    for b in bullets:
        p_b = tf_body.add_paragraph()
        p_b.text = f"  •  {b}"
        p_b.font.size = Pt(15)
        p_b.font.color.rgb = RGBColor(203, 213, 225)
        p_b.space_after = Pt(10)
        
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide, prs)
    add_slide_header(slide, "The Problem: Limitations of Current Education")
    tf_body = create_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
    p_sub = tf_body.paragraphs[0]
    p_sub.text = "Why Traditional Classrooms and Simple Chatbots Fail:"
    p_sub.font.bold = True
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = RGBColor(239, 68, 68)
    p_sub.space_after = Pt(15)
    points = [
        ("Lack of Adaptive Evaluation", "Typical online platforms offer static multiple-choice questions without explaining the root concepts of mistakes or tailoring explanations to the user's specific answers."),
        ("No Active Memory Loop", "Students consume study materials passively. Reading textbook files does not translate to conceptual memory retention without spaced quizzing and customized summaries."),
        ("Generic Chatbots are Context-Blind", "Common chatbot setups lack a RAG (Retrieval-Augmented Generation) layer, answering with general web knowledge rather than referencing local, user-provided textbooks and course materials."),
        ("Absence of Weak Topic Diagnostics", "Progress trackers display flat scores but fail to run semantic analysis over failures to warn students about underlying conceptual gaps like Recursion or Complex Variables.")
    ]
    for title, desc in points:
        p_title = tf_body.add_paragraph()
        p_title.text = f"✦  {title}"
        p_title.font.bold = True
        p_title.font.size = Pt(16)
        p_title.font.color.rgb = RGBColor(226, 232, 240)
        p_desc = tf_body.add_paragraph()
        p_desc.text = f"    {desc}"
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = RGBColor(148, 163, 184)
        p_desc.space_after = Pt(12)
        
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide, prs)
    add_slide_header(slide, "Proposed AI Solution: EduAI")
    tf_body = create_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
    p_intro = tf_body.paragraphs[0]
    p_intro.text = "EduAI acts as an end-to-end Personalized Academic Companion:"
    p_intro.font.bold = True
    p_intro.font.size = Pt(20)
    p_intro.font.color.rgb = RGBColor(167, 139, 250)
    p_intro.space_after = Pt(15)
    sol_points = [
        ("Retrieval-Augmented summaries", "Extracts text from PDFs/TXTs and segments it into vector blocks, enabling grounded study."),
        ("Interactive Adaptive Quizzing", "Generates tests (MCQs, Short-Answer, True/False) based on text difficulty."),
        ("Intelligent Evaluator Engine", "Runs answer grading, scores answers out of 5.0, gives feedback, and logs weak concepts."),
        ("Dynamic Study Planner", "Creates schedules, recommended revision topics, and master timelines based on student data."),
        ("Context-Aware Chat Tutor", "Enables chat queries with specific style controls like 'Explain in simple terms' or 'Explain like I'm 10'.")
    ]
    for title, desc in sol_points:
        p_t = tf_body.add_paragraph()
        p_t.text = f"✔  {title}:"
        p_t.font.bold = True
        p_t.font.size = Pt(16)
        p_t.font.color.rgb = RGBColor(255, 255, 255)
        p_d = tf_body.add_paragraph()
        p_d.text = f"    {desc}"
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = RGBColor(203, 213, 225)
        p_d.space_after = Pt(10)
        
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide, prs)
    add_slide_header(slide, "System Architecture")
    tf_body = create_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
    p_arch = tf_body.paragraphs[0]
    p_arch.text = "The architectural data flow of the application:"
    p_arch.font.bold = True
    p_arch.font.size = Pt(18)
    p_arch.font.color.rgb = RGBColor(96, 165, 250)
    p_arch.space_after = Pt(15)
    arch_steps = [
        ("1. Input Layer", "Student uploads local documents (PDF, TXT) and specifies difficulty settings."),
        ("2. Processing & Storage Layer", "PyPDF2 extracts text. The vectorizer chunks text and gets embeddings. SQLite stores progress logs."),
        ("3. LLM Agent Engine (Gemini)", "Receives prompts augmented with vector retrieval context to perform notes generation and evaluations."),
        ("4. Tracking & Execution Layer", "Calculates quiz errors, populates weak topics, and triggers a study plan generator."),
        ("5. User Interface (Streamlit)", "Provides a modern dashboard with Plotly visual aids.")
    ]
    for title, desc in arch_steps:
        p_s = tf_body.add_paragraph()
        p_s.text = f"•  {title}"
        p_s.font.bold = True
        p_s.font.size = Pt(15)
        p_s.font.color.rgb = RGBColor(226, 232, 240)
        p_sd = tf_body.add_paragraph()
        p_sd.text = f"    {desc}"
        p_sd.font.size = Pt(13)
        p_sd.font.color.rgb = RGBColor(148, 163, 184)
        p_sd.space_after = Pt(8)
        
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide, prs)
    add_slide_header(slide, "Technology Stack")
    tf_body = create_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
    stack = [
        ("Frontend & Dashboard", "Streamlit (Modern reactive UI, multi-page layout, custom styling)."),
        ("LLM Orchestrator", "Google Gemini API (gemini-pro for summarization, evaluations, and tutor chat)."),
        ("Retrieval Vector Store", "Gemini Embeddings Model + NumPy Cosine-Similarity for local semantic lookups."),
        ("Data Extraction", "PyPDF2 & pdfplumber for processing multi-page textbook PDF files."),
        ("Database Engine", "SQLite for lightweight local tracking of quiz configurations, grades, and weak topics."),
        ("Visual & Output Assets", "Plotly & Pandas for analytical graphs; python-pptx and ReportLab for automated assets.")
    ]
    for component, tech in stack:
        p_c = tf_body.add_paragraph()
        p_c.text = f"✦  {component}:"
        p_c.font.bold = True
        p_c.font.size = Pt(16)
        p_c.font.color.rgb = RGBColor(167, 139, 250)
        p_t = tf_body.add_paragraph()
        p_t.text = f"    {tech}"
        p_t.font.size = Pt(14)
        p_t.font.color.rgb = RGBColor(203, 213, 225)
        p_t.space_after = Pt(10)
        
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide, prs)
    add_slide_header(slide, "Detailed AI Workflow")
    tf_body = create_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
    workflow = [
        ("Step 1: Ingestion", "PDF text is parsed and split into 800-character chunks with a 150-character overlap."),
        ("Step 2: Vector Search", "Text chunks are indexed using `models/text-embedding-004` to yield relevant context."),
        ("Step 3: Adaptive Testing", "The quiz engine uses context to write custom MCQs and short answers."),
        ("Step 4: AI Assessment", "Student answers are evaluated by Gemini to score, explain errors, and tag concepts."),
        ("Step 5: Diagnostics & Action", "Incorrect concepts are logged, and the study plan generator structures a calendar.")
    ]
    for title, desc in workflow:
        p_wt = tf_body.add_paragraph()
        p_wt.text = f"⚙  {title}"
        p_wt.font.bold = True
        p_wt.font.size = Pt(15)
        p_wt.font.color.rgb = RGBColor(255, 255, 255)
        p_wd = tf_body.add_paragraph()
        p_wd.text = f"    {desc}"
        p_wd.font.size = Pt(13)
        p_wd.font.color.rgb = RGBColor(148, 163, 184)
        p_wd.space_after = Pt(8)
        
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide, prs)
    add_slide_header(slide, "EduAI Core Features")
    tf_body = create_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
    features = [
        ("Study Material Upload", "Streamlit uploader extracts text instantly and visualizes summaries."),
        ("AI Summarizer", "Presents summaries, key takeaway bullet lists, glosarry terms, and analogies."),
        ("Adaptive Quiz Generator", "Lets users customize test format (MCQ, Short answer, True/False) and set difficulty."),
        ("AI Evaluation & Critiques", "Provides grading marks, points out specific errors, and identifies weak concepts.")
    ]
    for title, desc in features:
        p_ft = tf_body.add_paragraph()
        p_ft.text = f"★  {title}:"
        p_ft.font.bold = True
        p_ft.font.size = Pt(16)
        p_ft.font.color.rgb = RGBColor(96, 165, 250)
        p_fd = tf_body.add_paragraph()
        p_fd.text = f"    {desc}"
        p_fd.font.size = Pt(14)
        p_fd.font.color.rgb = RGBColor(226, 232, 240)
        p_fd.space_after = Pt(10)
        
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide, prs)
    add_slide_header(slide, "Context-Aware AI Tutor")
    tf_body = create_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
    p_tutor_intro = tf_body.paragraphs[0]
    p_tutor_intro.text = "How the conversational RAG tutor guides learning:"
    p_tutor_intro.font.bold = True
    p_tutor_intro.font.size = Pt(18)
    p_tutor_intro.font.color.rgb = RGBColor(167, 139, 250)
    p_tutor_intro.space_after = Pt(15)
    details = [
        ("Grounded in Uploaded Content", "The tutor answers questions by matching them with textbook segments, preventing hallucinations."),
        ("Adaptive Personas", "Students can click standard persona templates to alter the tutor's vocabulary:"),
        ("  • 'Simple Words'", "Swaps complex technical definitions for clear explanations and analogies."),
        ("  • 'Explain like I'm 10'", "Uses simple examples appropriate for a younger audience."),
        ("  • 'Give another example'", "Generates custom real-world examples to explain theoretical concepts.")
    ]
    for title, desc in details:
        p_dt = tf_body.add_paragraph()
        p_dt.text = f"  {title}"
        if title.startswith("  •"):
            p_dt.font.bold = False
            p_dt.font.size = Pt(14)
            p_dt.font.color.rgb = RGBColor(148, 163, 184)
        else:
            p_dt.font.bold = True
            p_dt.font.size = Pt(15)
            p_dt.font.color.rgb = RGBColor(226, 232, 240)
        if desc:
            p_dd = tf_body.add_paragraph()
            p_dd.text = f"    {desc}"
            p_dd.font.size = Pt(13)
            p_dd.font.color.rgb = RGBColor(203, 213, 225)
            p_dd.space_after = Pt(8)
            
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide, prs)
    add_slide_header(slide, "SDG 4 Quality Education Impact")
    tf_body = create_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
    impacts = [
        ("Democratic Access to Tutoring", "High-quality academic assistance is made accessible to all students with a web browser, bypassing private tutoring costs."),
        ("Personalized Learning Pathways", "Translates student performance data into targeted, customized revision tasks to improve concept comprehension."),
        ("Mitigating Classroom Gaps", "Identifies conceptual weaknesses before final exams, helping teachers and students address learning gaps."),
        ("Enabling Independent Learning", "Supports student agency by providing a safe space to ask clarifying questions and practice quizzes.")
    ]
    for title, desc in impacts:
        p_it = tf_body.add_paragraph()
        p_it.text = f"✓  {title}"
        p_it.font.bold = True
        p_it.font.size = Pt(16)
        p_it.font.color.rgb = RGBColor(96, 165, 250)
        p_id = tf_body.add_paragraph()
        p_id.text = f"    {desc}"
        p_id.font.size = Pt(14)
        p_id.font.color.rgb = RGBColor(203, 213, 225)
        p_id.space_after = Pt(10)
        
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide, prs)
    add_slide_header(slide, "Future Scope & Extensions")
    tf_body = create_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
    scopes = [
        ("Multi-lingual Support", "Translate textbook summaries and conduct tutor sessions in native regional languages to broaden accessibility."),
        ("LMS API Integration", "Sync student weak topics directly with platforms like Canvas, Moodle, or Google Classroom for teacher review."),
        ("OCR for Handwritten Notes", "Enable scanning and uploading handwritten lecture sheets or diagram files for RAG processing."),
        ("Offline Embedding Sync", "Implement lightweight offline models to run on-device, resolving internet dependency in remote areas.")
    ]
    for title, desc in scopes:
        p_st = tf_body.add_paragraph()
        p_st.text = f"✦  {title}"
        p_st.font.bold = True
        p_st.font.size = Pt(16)
        p_st.font.color.rgb = RGBColor(167, 139, 250)
        p_sd = tf_body.add_paragraph()
        p_sd.text = f"    {desc}"
        p_sd.font.size = Pt(14)
        p_sd.font.color.rgb = RGBColor(203, 213, 225)
        p_sd.space_after = Pt(10)
        
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide, prs)
    add_slide_header(slide, "Conclusion")
    tf_body = create_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
    p_c1 = tf_body.paragraphs[0]
    p_c1.text = "EduAI demonstrates that AI can do more than answer queries; it can serve as a personalized, context-aware study companion."
    p_c1.font.size = Pt(18)
    p_c1.font.color.rgb = RGBColor(255, 255, 255)
    p_c1.space_after = Pt(15)
    p_c2 = tf_body.add_paragraph()
    p_c2.text = "Key Takeaways:"
    p_c2.font.bold = True
    p_c2.font.size = Pt(18)
    p_c2.font.color.rgb = RGBColor(96, 165, 250)
    p_c2.space_after = Pt(10)
    takeaways = [
        "Aligns AI capabilities with UN SDG 4 to support inclusive, high-quality learning.",
        "Uses SQLite records to transition from context-blind chat to structured tracking.",
        "Provides a scalable, lightweight architecture suitable for deployment on low-cost devices."
    ]
    for t in takeaways:
        p_t = tf_body.add_paragraph()
        p_t.text = f"  ✔  {t}"
        p_t.font.size = Pt(16)
        p_t.font.color.rgb = RGBColor(203, 213, 225)
        p_t.space_after = Pt(10)
    p_q = tf_body.add_paragraph()
    p_q.text = "Thank You! Questions & Discussion."
    p_q.font.bold = True
    p_q.font.size = Pt(20)
    p_q.font.color.rgb = RGBColor(167, 139, 250)
    p_q.space_before = Pt(30)
    p_q.alignment = PP_ALIGN.CENTER
    
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)

if __name__ == "__main__":
    generate_presentation()